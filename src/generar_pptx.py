"""
Genera la presentación PPTX completa del Datathon FinanCrece S.A. — ESAN 2026
Equipo: Armonitech
Slides: 15 detalladas con todas las preguntas del jurado respondidas
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy
from pathlib import Path
from PIL import Image
import io

FIG = Path("notebooks/figures")
OUT = Path("notebooks/GRUPOArmonitech_MEJORADO.pptx")

# ── Paleta de colores ──────────────────────────────────────────
BG_DARK   = RGBColor(0x0A, 0x16, 0x28)   # fondo oscuro
BG_CARD   = RGBColor(0x0D, 0x1F, 0x38)   # card azul marino
ACCENT    = RGBColor(0x00, 0xE5, 0xFF)   # cyan eléctrico
GOLD      = RGBColor(0xF5, 0xA6, 0x23)   # naranja dorado
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xA0, 0xAE, 0xC0)   # gris claro
GREEN     = RGBColor(0x22, 0xC5, 0x5E)
RED       = RGBColor(0xEF, 0x44, 0x44)
AMBER     = RGBColor(0xF5, 0x9E, 0x0B)
BLUE      = RGBColor(0x2D, 0x6D, 0xB5)
DARKBLUE2 = RGBColor(0x1E, 0x3A, 0x5F)

prs = Presentation()
prs.slide_width  = Inches(16)
prs.slide_height = Inches(9)

BLANK = prs.slide_layouts[6]   # layout en blanco

# ── Helpers ────────────────────────────────────────────────────
def slide_bg(slide, color=BG_DARK):
    """Pinta el fondo del slide"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_text_lines(slide, lines, l, t, w, h,
                   size=14, color=WHITE, line_spacing_pt=None):
    """
    lines: list of (text, bold, color_override_or_None)
    """
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, bld, col) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(size)
        run.font.bold = bld
        run.font.color.rgb = col if col else color

def add_image_safe(slide, img_path, l, t, w, h=None):
    path = Path(img_path)
    if not path.exists():
        return
    if h is None:
        slide.shapes.add_picture(str(path), Inches(l), Inches(t), width=Inches(w))
    else:
        slide.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w), Inches(h))

def header_bar(slide, title, subtitle=None, title_size=30, sub_size=16):
    """Barra de título en la parte superior"""
    add_rect(slide, 0, 0, 16, 1.45, BG_CARD)
    add_rect(slide, 0, 1.35, 16, 0.06, ACCENT)   # línea de acento
    add_text(slide, title, 0.35, 0.10, 13, 0.8,
             size=title_size, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.82, 13, 0.52,
                 size=sub_size, color=LIGHT, align=PP_ALIGN.LEFT)

def footer(slide, num, total=15):
    add_rect(slide, 0, 8.72, 16, 0.28, DARKBLUE2)
    add_text(slide, f"FinanCrece S.A. — Datathon ESAN 2026  |  Equipo Armonitech",
             0.2, 8.74, 12, 0.22, size=8, color=LIGHT)
    add_text(slide, f"{num} / {total}", 15.0, 8.74, 0.8, 0.22,
             size=8, color=LIGHT, align=PP_ALIGN.RIGHT)

def tag(slide, txt, l, t, w=1.6, h=0.30, bg=ACCENT, fg=BG_DARK):
    add_rect(slide, l, t, w, h, bg)
    add_text(slide, txt, l+0.05, t+0.02, w-0.10, h-0.04,
             size=9, bold=True, color=fg, align=PP_ALIGN.CENTER)

def kpi_card(slide, l, t, w, h, value, label, color=ACCENT):
    add_rect(slide, l, t, w, h, BG_CARD)
    add_rect(slide, l, t, w, 0.04, color)  # top accent line
    add_text(slide, value, l+0.08, t+0.12, w-0.16, h*0.48,
             size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, l+0.08, t+h*0.55, w-0.16, h*0.40,
             size=10, color=LIGHT, align=PP_ALIGN.CENTER)

def table_simple(slide, headers, rows, l, t, w, h,
                 header_bg=DARKBLUE2, row_bg=BG_CARD, alt_bg=BG_DARK):
    n_cols = len(headers)
    n_rows = len(rows)
    col_w = w / n_cols
    row_h = h / (n_rows + 1)

    # Header
    for ci, hdr in enumerate(headers):
        add_rect(slide, l + ci*col_w, t, col_w - 0.02, row_h - 0.02, header_bg)
        add_text(slide, hdr, l + ci*col_w + 0.05, t + 0.02,
                 col_w - 0.12, row_h - 0.06, size=9, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Rows
    for ri, row in enumerate(rows):
        bg = row_bg if ri % 2 == 0 else alt_bg
        for ci, cell in enumerate(row):
            add_rect(slide, l + ci*col_w, t + (ri+1)*row_h, col_w - 0.02, row_h - 0.02, bg)
            col_c = WHITE
            if str(cell).startswith('✅'): col_c = GREEN
            elif str(cell).startswith('🔴'): col_c = RED
            elif str(cell).startswith('⚠️'): col_c = AMBER
            add_text(slide, str(cell),
                     l + ci*col_w + 0.05, t + (ri+1)*row_h + 0.02,
                     col_w - 0.12, row_h - 0.06,
                     size=8.5, color=col_c, align=PP_ALIGN.LEFT)

def bullet_block(slide, items, l, t, w, h, size=12, bullet="▶", color=WHITE, accent_col=None):
    """items: list of (text, is_bold, color_override)"""
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, str):
            txt, bld, col = item, False, color
        else:
            txt, bld, col = item[0], item[1], item[2] if len(item)>2 else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"{bullet} {txt}"
        r.font.size = Pt(size)
        r.font.bold = bld
        r.font.color.rgb = col

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
# Degradado visual: rectángulos de fondo
add_rect(sl, 0, 0, 8, 9, RGBColor(0x0A,0x16,0x28))
add_rect(sl, 8, 0, 8, 9, RGBColor(0x0D,0x1F,0x38))
add_rect(sl, 7.8, 0, 0.06, 9, ACCENT)

# Lado izquierdo
add_text(sl, "DATATHON ESAN 2026", 0.5, 0.8, 7, 0.5,
         size=13, bold=False, color=ACCENT, italic=True)
add_text(sl, "FinanCrece S.A.", 0.5, 1.35, 7.2, 1.1,
         size=52, bold=True, color=WHITE)
add_text(sl, "Sistema de Scoring de\nRiesgo Crediticio", 0.5, 2.6, 7, 1.2,
         size=28, bold=False, color=GOLD)
add_text(sl,
    "Predecir el riesgo antes del desembolso:\n"
    "la diferencia entre rentabilidad y mora.",
    0.5, 4.1, 7, 0.9, size=15, color=LIGHT, italic=True)

# Separador
add_rect(sl, 0.5, 5.2, 5.5, 0.04, ACCENT)

add_text(sl, "EQUIPO ARMONITECH", 0.5, 5.45, 7, 0.4,
         size=14, bold=True, color=ACCENT)
add_text(sl, "Mayo 2026  ·  Lima, Perú", 0.5, 5.9, 7, 0.4,
         size=13, color=LIGHT)

# Lado derecho — métricas clave
kpi_card(sl, 8.5,  1.0, 3.3, 1.4, "AUC = 0.833", "ROC-AUC del modelo campeón", ACCENT)
kpi_card(sl, 12.1, 1.0, 3.3, 1.4, "Gini = 0.665", "Coeficiente de discriminación", GOLD)
kpi_card(sl, 8.5,  2.65, 3.3, 1.4, "KS = 0.633", "Separación pagadores/morosos", GREEN)
kpi_card(sl, 12.1, 2.65, 3.3, 1.4, "+107% ROI", "vs. aprobar todos [SUPUESTO]", AMBER)
kpi_card(sl, 8.5,  4.3,  3.3, 1.4, "800 clientes", "Dataset de entrenamiento", BLUE)
kpi_card(sl, 12.1, 4.3,  3.3, 1.4, "62 features", "Variables diseñadas", BLUE)

add_rect(sl, 8.5, 6.0, 6.9, 2.65, RGBColor(0x05,0x0E,0x1C))
add_text(sl, "Modelo campeón: LightGBM regularizado\n"
             "Pipeline anti-leakage · Split estratificado 60/20/20\n"
             "Política de 3 bandas · Interpretabilidad SHAP\n"
             "Optimización de threshold por ROI financiero",
         8.65, 6.15, 6.6, 2.3, size=13, color=LIGHT)
footer(sl, 1)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — EL PROBLEMA: ¿POR QUÉ EL CRÉDITO FALLA?
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "El Problema: ¿Por qué el crédito falla?",
           "P1 & P2 · Causas del incumplimiento en FinanCrece S.A.")

# Contexto
add_rect(sl, 0.3, 1.6, 7.0, 3.7, BG_CARD)
add_rect(sl, 0.3, 1.6, 0.07, 3.7, RED)
add_text(sl, "La mora subió de 4.2% → 7.8%", 0.55, 1.7, 6.6, 0.55,
         size=18, bold=True, color=RED)
add_text(sl, "+3.6 puntos porcentuales en 18 meses", 0.55, 2.2, 6.6, 0.4,
         size=13, color=LIGHT)

bullet_block(sl, [
    ("Sin un modelo de scoring, el banco aprueba subjetivamente", False, LIGHT),
    ("Los analistas no pueden procesar 800+ solicitudes con igual rigor", False, LIGHT),
    ("La mora destruye rentabilidad y exige provisiones de capital adicionales", False, LIGHT),
    ("El costo de un default (-$3,000) supera 6.7x la ganancia de un buen crédito (+$450)", True, GOLD),
], 0.55, 2.65, 6.6, 2.4, size=12)

# Señales detectadas
add_text(sl, "Señales de riesgo detectadas en los datos:", 7.6, 1.6, 8.0, 0.45,
         size=14, bold=True, color=ACCENT)
table_simple(sl,
    ["Variable", "Grupo", "Default Rate"],
    [
        ["checking_status", "Cuenta negativa (<0)", "🔴 47.4%"],
        ["credit_history", "Sin historial previo", "🔴 65.6%"],
        ["employment", "Empleo < 1 año", "🔴 40.4%"],
        ["savings_status", "Ahorros < 100 USD", "🔴 35.3%"],
        ["housing", "Vivienda alquilada", "🔴 39.0%"],
        ["duration", "Plazo > 24 meses", "⚠️ +12pp"],
        ["checking_status", "Sin cuenta corriente", "✅ 10.9%"],
        ["employment", "4 a 7 años empleo", "✅ 21.4%"],
    ],
    7.6, 2.1, 8.0, 4.0
)

add_text(sl, "Promedio de default del portafolio: 29.5% (236 / 800 clientes)",
         0.3, 5.45, 15.4, 0.4, size=11, color=LIGHT, italic=True)
add_image_safe(sl, FIG/"eda_target_y_checking.png", 0.3, 5.85, 15.4, 2.6)
footer(sl, 2)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — PROPÓSITO DEL CRÉDITO
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "¿En qué se usa el crédito y quién incumple más?",
           "P3 · Propósito del crédito vs. probabilidad de default")

add_image_safe(sl, FIG/"ppt_proposito_vs_default.png", 0.3, 1.55, 9.8, 5.1)

add_rect(sl, 10.3, 1.55, 5.4, 5.1, BG_CARD)
add_rect(sl, 10.3, 1.55, 0.06, 5.1, GOLD)
add_text(sl, "Hallazgos clave", 10.5, 1.65, 5.0, 0.45,
         size=14, bold=True, color=GOLD)

table_simple(sl,
    ["Propósito", "DR", "Acción"],
    [
        ["Radio/TV", "21.8%", "✅ Aprobar"],
        ["Auto Usado", "15.4%", "✅ Aprobar"],
        ["Muebles/Equipo", "31.3%", "⚠️ Evaluar"],
        ["Auto Nuevo", "36.6%", "⚠️ Evaluar"],
        ["Negocio", "38.2%", "🔴 Reforzar"],
        ["Reparaciones", "38.9%", "🔴 Reforzar"],
        ["Educación", "40.5%", "🔴 Garantías"],
    ],
    10.3, 2.2, 5.4, 3.5
)
add_text(sl,
    "El propósito del crédito es una\n"
    "variable moderadora del riesgo:\n"
    "bienes duraderos < riesgo que\n"
    "inversión productiva o educación.",
    10.35, 5.75, 5.2, 0.85, size=11, color=LIGHT, italic=True)

add_text(sl,
    "Pregunta del jurado: ¿Qué hace que las personas soliciten un crédito?\n"
    "→ Consumo (Radio/TV 28%), Autos (32%), Negocio (10%), Educación (5%)",
    0.3, 6.75, 15.4, 0.55, size=10, color=ACCENT, italic=True)
footer(sl, 3)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — EMPLEO: TIEMPO IDEAL
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "¿Cuál es el tiempo de empleo ideal para un crédito?",
           "P4 · Employment vs. Default Rate — Datos reales del portafolio")

add_image_safe(sl, FIG/"ppt_empleo_tiempo_vs_default.png", 0.3, 1.55, 10.2, 5.0)

add_rect(sl, 10.8, 1.55, 4.9, 5.0, BG_CARD)
add_rect(sl, 10.8, 1.55, 0.06, 5.0, GREEN)
add_text(sl, "Interpretación bancaria", 11.0, 1.65, 4.5, 0.45,
         size=13, bold=True, color=GREEN)

bullet_block(sl, [
    ("4–7 años: ÓPTIMO → DR = 21.4% ✅", True, GREEN),
    ("7+ años: MUY BUENO → DR = 22.4% ✅", True, GREEN),
    ("1–4 años: MODERADO → DR = 33.1% ⚠️", False, AMBER),
    ("< 1 año: ALTO RIESGO → DR = 40.4% 🔴", False, RED),
    ("Desempleado: DR = 30.2% (n pequeño)", False, LIGHT),
    ("", False, WHITE),
    ("REGLA BANCARIA SUGERIDA:", True, GOLD),
    ("Exigir mínimo 1 año de empleo continuo para créditos > $2,000", False, LIGHT),
    ("Para créditos > $5,000: mínimo 4 años de empleo", False, LIGHT),
    ("", False, WHITE),
    ("El modelo penaliza empleo < 1 año con employment_risk_ordinal = 3 (máximo riesgo)", False, LIGHT),
], 11.0, 2.2, 4.5, 4.1, size=11)

add_text(sl,
    "→ El tiempo en el empleo actual es un proxy de estabilidad de ingresos y compromiso de pago",
    0.3, 6.7, 15.4, 0.45, size=11, color=ACCENT, italic=True)
footer(sl, 4)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — ESTADO CIVIL Y VIVIENDA
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Estado Civil y Vivienda: ¿Influyen en el pago?",
           "P5, P6, P7, P8 · Personal status y housing vs. Default Rate")

# Imagen estado civil
add_text(sl, "Estado Civil", 0.3, 1.6, 7.5, 0.38,
         size=14, bold=True, color=GOLD)
add_image_safe(sl, FIG/"ppt_estado_civil_vs_default.png", 0.3, 2.0, 7.5, 3.8)

# Imagen vivienda
add_text(sl, "Tipo de Vivienda", 8.2, 1.6, 7.5, 0.38,
         size=14, bold=True, color=GOLD)
add_image_safe(sl, FIG/"ppt_vivienda_vs_default.png", 8.2, 2.0, 7.5, 3.8)

# Conclusiones en el footer del slide
add_rect(sl, 0.3, 5.9, 7.5, 2.6, BG_CARD)
add_rect(sl, 0.3, 5.9, 0.06, 2.6, BLUE)
bullet_block(sl, [
    ("✅ Hombres casados/viudos: menor DR (24.7%)", True, GREEN),
    ("✅ Hombres solteros: DR = 25.8%", True, GREEN),
    ("🔴 Mujer divorciada/dependiente: DR = 36.1%", True, RED),
    ("🔴 Hombre divorciado: DR = 36.4%", True, RED),
    ("→ Diferencia de 11pp entre mejor y peor grupo", False, ACCENT),
    ("→ Estado civil = señal de estabilidad del hogar", False, LIGHT),
], 0.5, 6.0, 7.1, 2.3, size=11)

add_rect(sl, 8.2, 5.9, 7.5, 2.6, BG_CARD)
add_rect(sl, 8.2, 5.9, 0.06, 2.6, BLUE)
bullet_block(sl, [
    ("✅ Vivienda propia: DR = 25.6% (71% del portafolio)", True, GREEN),
    ("🔴 Alquilada: DR = 39.0% → 15pp más riesgo", True, RED),
    ("🔴 Cedida (for free): DR = 39.8%", True, RED),
    ("→ Vivienda propia = colateral implícito de estabilidad", False, ACCENT),
    ("→ Puede usarse como criterio de elegibilidad mínima", False, LIGHT),
    ("→ Alquilada puede exigir co-deudor o garantía", False, LIGHT),
], 8.4, 6.0, 7.1, 2.3, size=11)
footer(sl, 5)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — EDA: CUENTA CORRIENTE (VARIABLE REINA)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "La variable más predictiva: Estado de Cuenta Corriente",
           "P11 · checking_status — El discriminador más potente del portafolio")

add_text(sl, "Una sola variable discrimina 5.6x entre el mejor y el peor grupo:", 0.3, 1.6, 15.4, 0.4,
         size=13, color=LIGHT)

table_simple(sl,
    ["checking_status", "Clientes (n)", "Default Rate", "Diferencia vs promedio", "Acción recomendada"],
    [
        ["Sin cuenta corriente (no checking)", "274 (34%)", "✅ 10.9%", "-18.6pp", "✅ Aprobar con normalidad"],
        ["Cuenta >= 200 DM (buena)", "63 (8%)", "✅ 21.8%", "-7.7pp", "✅ Aprobar — bajo riesgo"],
        ["Cuenta 0-200 DM (baja)", "269 (34%)", "⚠️ 41.4%", "+11.9pp", "⚠️ Condicionar crédito"],
        ["Cuenta negativa (< 0)", "194 (24%)", "🔴 47.4%", "+17.9pp", "🔴 Rechazar o garantías"],
    ],
    0.3, 2.1, 15.4, 3.0
)

add_rect(sl, 0.3, 5.25, 7.4, 3.2, BG_CARD)
add_rect(sl, 0.3, 5.25, 0.06, 3.2, ACCENT)
add_text(sl, "Interpretación bancaria", 0.55, 5.35, 6.8, 0.38,
         size=13, bold=True, color=ACCENT)
bullet_block(sl, [
    ("Cuenta negativa = está gastando más de lo que tiene", False, LIGHT),
    ("Señal de stress financiero activo al momento de la solicitud", False, LIGHT),
    ("El banco debería EXIGIR estado de cuenta como requisito mínimo", True, GOLD),
    ("Sin esta variable, el riesgo del portafolio se subestima en ~18pp", True, RED),
], 0.55, 5.8, 6.9, 2.4, size=11)

add_rect(sl, 8.0, 5.25, 7.7, 3.2, BG_CARD)
add_rect(sl, 8.0, 5.25, 0.06, 3.2, GOLD)
add_text(sl, "Encodificación para el modelo", 8.25, 5.35, 7.2, 0.38,
         size=13, bold=True, color=GOLD)
table_simple(sl,
    ["Valor original", "checking_risk_ordinal", "Peso en score"],
    [
        ["no checking", "0", "Sin penalización"],
        [">= 200", "1", "Baja penalización"],
        ["0 <= X < 200", "2", "Penalización media"],
        ["< 0", "3", "Penalización alta"],
    ],
    8.0, 5.75, 7.7, 2.5
)
footer(sl, 6)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — FEATURE ENGINEERING
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Feature Engineering Bancario: 20 → 62 variables",
           "Transformaciones fila a fila — Anti-leakage garantizado")

# Left panel
add_rect(sl, 0.3, 1.6, 7.5, 6.8, BG_CARD)
add_rect(sl, 0.3, 1.6, 0.06, 6.8, ACCENT)
add_text(sl, "Grupos de features creadas", 0.55, 1.7, 7.0, 0.4,
         size=13, bold=True, color=ACCENT)

grupos = [
    ("CRÉDITO", "duration_largo/corto · log_duration · log_credit_amount · monto_alto · monto_bajo", BLUE),
    ("CARGA FINANCIERA", "carga_financiera = amount × duration · cuota_estimada = amount ÷ duration · log_carga", BLUE),
    ("CUENTA CORRIENTE", "checking_risk_ordinal (0–3) · cuenta_negativa · sin_cuenta · cuenta_buena", ACCENT),
    ("HISTORIAL", "history_risk_ordinal (1–4) · historial_limpio · historial_critico · historial_delay", ACCENT),
    ("AHORROS", "savings_risk_ordinal (0–4) · sin_ahorros · ahorros_altos", GOLD),
    ("EMPLEO", "employment_risk_ordinal (0–4) · desempleado · empleo_estable", GOLD),
    ("INTERACCIONES", "riesgo_combinado = check × hist · negativo_y_monto_alto · sin_reservas · carga_vs_cuenta", GREEN),
    ("SCORE COMPUESTO", "30×check + 20×hist + 15×ahorro + 10×empleo + 5×log_carga", AMBER),
]
y0 = 2.2
for nom, desc, col in grupos:
    add_rect(sl, 0.35, y0, 0.08, 0.55, col)
    add_text(sl, nom, 0.55, y0, 2.5, 0.25, size=9, bold=True, color=col)
    add_text(sl, desc, 0.55, y0+0.26, 7.0, 0.28, size=8.5, color=LIGHT)
    y0 += 0.72

# Right panel
add_rect(sl, 8.1, 1.6, 7.6, 3.8, BG_CARD)
add_rect(sl, 8.1, 1.6, 0.06, 3.8, GOLD)
add_text(sl, "Top 10 Variables por Importancia (LightGBM)", 8.35, 1.7, 7.1, 0.4,
         size=13, bold=True, color=GOLD)
add_image_safe(sl, FIG/"ppt_feature_importance_jurado.png", 8.1, 2.1, 7.6, 3.2)

add_rect(sl, 8.1, 5.55, 7.6, 2.85, BG_CARD)
add_rect(sl, 8.1, 5.55, 0.06, 2.85, GREEN)
add_text(sl, "Regla Anti-Leakage", 8.35, 5.65, 7.1, 0.38,
         size=13, bold=True, color=GREEN)
bullet_block(sl, [
    ("feature_builder.py: transformaciones FILA A FILA — sin fit", True, ACCENT),
    ("El preprocessor (imputador + scaler) se ajusta SOLO en train", False, LIGHT),
    ("Se aplica idéntico a val y test → sin contaminación", False, LIGHT),
    ("Split: 60% train | 20% val | 20% test (stratified)", True, GOLD),
    ("0 variables sospechosas de leakage detectadas", True, GREEN),
], 8.35, 6.1, 7.1, 2.1, size=11)
footer(sl, 7)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — MODELOS Y MÉTRICAS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Comparativa de Modelos: 5 candidatos, 1 campeón",
           "P11 · Selección del modelo campeón por AUC, Gini, KS y estabilidad")

table_simple(sl,
    ["Modelo", "AUC Val", "Gini Val", "KS Val", "Brier", "Lift@10%", "Gap Overfit", "Veredicto"],
    [
        ["Dummy Baseline",        "0.468", "-0.063", "0.125", "0.444", "1.06x", "0.064", "—"],
        ["Logistic Regression",   "0.783", "0.565",  "0.479", "0.182", "2.55x", "0.095", "⚠️ Bueno"],
        ["Random Forest",         "0.812", "0.623",  "0.540", "0.157", "2.98x", "0.147", "⚠️ Bueno"],
        ["HistGradientBoosting",  "0.794", "0.589",  "0.491", "0.181", "2.13x", "0.206", "🔴 Overfit"],
        ["XGBoost",               "0.818", "0.635",  "0.532", "0.160", "2.55x", "0.179", "⚠️ Bueno"],
        ["CatBoost",              "0.813", "0.626",  "0.508", "0.159", "2.98x", "0.172", "⚠️ Bueno"],
        ["LightGBM Regularizado", "0.833", "0.665",  "0.633", "0.157", "2.34x", "0.094", "✅ CAMPEÓN"],
    ],
    0.3, 1.6, 15.4, 4.0
)

# KPIs finales
kpi_card(sl, 0.3,  5.8, 3.5, 1.5, "0.833",  "ROC-AUC (Validación)", ACCENT)
kpi_card(sl, 4.1,  5.8, 3.5, 1.5, "0.665",  "Gini (Discriminación)", GOLD)
kpi_card(sl, 7.9,  5.8, 3.5, 1.5, "0.633",  "KS (Separación colas)", GREEN)
kpi_card(sl, 11.7, 5.8, 3.9, 1.5, "0.157",  "Brier (Calibración)", BLUE)

add_text(sl,
    "Criterio de selección: Mayor AUC + menor Gap de overfit + mayor KS.\n"
    "LightGBM regularizado (learning_rate=0.05, max_depth=4, reg_alpha=0.5, reg_lambda=1.0) logra el mejor balance.",
    0.3, 7.45, 15.4, 0.6, size=10, color=LIGHT, italic=True)
footer(sl, 8)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — CURVAS DE EVALUACIÓN
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Panel de Evaluación del Modelo Campeón",
           "ROC · Lift · Distribución de scores · Feature Importance")

add_image_safe(sl, FIG/"panel_evaluacion_final.png", 0.3, 1.55, 15.4, 7.0)
footer(sl, 9)

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — MONTO, PLAZO Y CUOTA
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Monto, Plazo y Cuota: La Tríada del Riesgo",
           "P14, P15 · credit_amount × duration × installment_commitment")

add_image_safe(sl, FIG/"ppt_monto_duracion_riesgo.png", 0.3, 1.55, 9.5, 3.8)
add_image_safe(sl, FIG/"ppt_cuota_comprometida.png",    0.3, 5.45, 9.5, 3.0)

add_rect(sl, 10.1, 1.55, 5.6, 3.8, BG_CARD)
add_rect(sl, 10.1, 1.55, 0.06, 3.8, GOLD)
add_text(sl, "Monto × Duración", 10.35, 1.65, 5.1, 0.4,
         size=13, bold=True, color=GOLD)
table_simple(sl,
    ["Carga Financiera", "Default Rate"],
    [
        ["Q1 — Baja",   "✅ 21.4%"],
        ["Q2",          "⚠️ 27.0%"],
        ["Q3",          "⚠️ 29.5%"],
        ["Q4 — Alta",   "🔴 39.2%"],
    ],
    10.1, 2.15, 5.6, 2.55
)
add_text(sl, "Plazos > 24 meses = señal de alerta automática",
         10.35, 4.65, 5.2, 0.5, size=10, color=AMBER, bold=True)

add_rect(sl, 10.1, 5.45, 5.6, 3.0, BG_CARD)
add_rect(sl, 10.1, 5.45, 0.06, 3.0, ACCENT)
add_text(sl, "Cuota Comprometida", 10.35, 5.55, 5.1, 0.4,
         size=13, bold=True, color=ACCENT)
table_simple(sl,
    ["Nivel cuota", "Default Rate"],
    [
        ["Nivel 1 (bajo%)", "✅ 23.4%"],
        ["Nivel 2",         "⚠️ 26.3%"],
        ["Nivel 3",         "⚠️ 29.8%"],
        ["Nivel 4 (alto%)", "🔴 32.8%"],
    ],
    10.1, 5.95, 5.6, 2.1
)
add_text(sl, "Relación CASI LINEAL: a mayor cuota → mayor mora",
         10.35, 8.15, 5.2, 0.3, size=10, color=AMBER, bold=True)
footer(sl, 10)

# ══════════════════════════════════════════════════════════════
# SLIDE 11 — AHORROS: ¿LINEAL O NO?
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "¿Los Ahorros Tienen Impacto Lineal en el Riesgo?",
           "P13 · savings_status — El umbral real de seguridad crediticia")

add_image_safe(sl, FIG/"ppt_ahorros_linealidad.png", 0.3, 1.55, 10.0, 4.5)

add_rect(sl, 10.5, 1.55, 5.2, 4.5, BG_CARD)
add_rect(sl, 10.5, 1.55, 0.06, 4.5, BLUE)
add_text(sl, "¿Es lineal el impacto?", 10.75, 1.65, 4.7, 0.4,
         size=13, bold=True, color=BLUE)
table_simple(sl,
    ["Nivel de ahorros", "Default Rate", "Insight"],
    [
        ["< 100 (sin ahorros)",   "🔴 35.3%", "Máximo riesgo"],
        ["100 – 500 (bajo)",      "🔴 35.1%", "Sin mejora apreciable"],
        ["500 – 1000 (medio)",    "✅ 19.6%", "CAÍDA ABRUPTA"],
        [">= 1000 (alto)",        "✅ 12.2%", "Mínimo riesgo"],
        ["Sin datos de ahorros",  "✅ 16.0%", "Proxy de patrimonio"],
    ],
    10.5, 2.1, 5.2, 3.5
)

add_rect(sl, 0.3, 6.2, 15.4, 2.2, BG_CARD)
add_rect(sl, 0.3, 6.2, 0.06, 2.2, AMBER)
add_text(sl, "Hallazgo no obvio — El umbral está en 500 USD", 0.55, 6.3, 8.0, 0.4,
         size=14, bold=True, color=AMBER)
bullet_block(sl, [
    ("Por debajo de 500 USD de ahorros, el riesgo NO mejora (35.3% vs 35.1% → prácticamente igual)", False, LIGHT),
    ("A partir de 500 USD se produce una caída abrupta de 15.5pp → frontera real de seguridad", True, ACCENT),
    ("Recomendación: usar savings >= 500 como criterio de elegibilidad para créditos sin garantía", True, GOLD),
    ("El banco podría ofrecer tasa preferencial a clientes con ahorros > 1,000 USD (DR = 12.2%)", False, LIGHT),
], 0.55, 6.75, 14.8, 1.5, size=11)
footer(sl, 11)

# ══════════════════════════════════════════════════════════════
# SLIDE 12 — RIESGOS OCULTOS E INCLUSIÓN FINANCIERA
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Riesgos Ocultos e Inclusión Financiera",
           "P9, P19, P20 · Clientes sin historial — El dilema de inclusión responsable")

add_image_safe(sl, FIG/"ppt_inclusion_riesgo_oculto.png", 0.3, 1.55, 9.8, 4.5)

add_rect(sl, 10.3, 1.55, 5.4, 2.1, BG_CARD)
add_rect(sl, 10.3, 1.55, 0.06, 2.1, RED)
add_text(sl, "El problema de inclusión", 10.55, 1.65, 5.0, 0.38,
         size=12, bold=True, color=RED)
bullet_block(sl, [
    ("32 clientes sin historial crediticio (4% del portafolio)", False, LIGHT),
    ("Tasa de default = 65.6% → el mayor riesgo del portafolio", True, RED),
    ("Excluirlos = discriminación financiera", False, LIGHT),
    ("Incluirlos sin score = destrucción de capital", True, RED),
], 10.55, 2.1, 5.1, 1.35, size=10)

add_rect(sl, 10.3, 3.75, 5.4, 2.3, BG_CARD)
add_rect(sl, 10.3, 3.75, 0.06, 2.3, GREEN)
add_text(sl, "La solución del modelo", 10.55, 3.85, 5.0, 0.38,
         size=12, bold=True, color=GREEN)
bullet_block(sl, [
    ("En vez de excluir → CUANTIFICAR el riesgo individual", True, ACCENT),
    ("Asignar a banda media → línea reducida al 50% [SUPUESTO]", False, LIGHT),
    ("Variables proxy compensan la falta de historial:", False, LIGHT),
    ("  Empleo estable + vivienda propia + edad madura", False, GOLD),
    ("Si paga bien → graduación a banda baja (inclusión responsable)", True, GREEN),
], 10.55, 4.3, 5.1, 1.6, size=10)

add_rect(sl, 0.3, 6.2, 15.4, 2.2, BG_CARD)
add_rect(sl, 0.3, 6.2, 0.06, 2.2, RED)
add_text(sl, "Señales de Riesgo Oculto en Clientes Aparentemente Solventes", 0.55, 6.3, 14.8, 0.38,
         size=13, bold=True, color=RED)
table_simple(sl,
    ["Señal oculta", "Variable proxy", "Default Rate", "Riesgo oculto"],
    [
        ["Sin historial formal",     "credit_history = no credits",    "🔴 65.6%", "Muy alto"],
        ["Cuenta en negativo",       "checking_status = <0",           "🔴 47.4%", "Alto"],
        ["Empleo muy reciente",      "employment = <1",                "🔴 40.4%", "Alto"],
        ["Propósito improductivo",   "purpose = education/repairs",    "⚠️ 40.5%", "Medio-Alto"],
        ["Cuota máxima comprometida","installment_commitment = 4",     "⚠️ 32.8%", "Medio"],
    ],
    0.3, 6.7, 15.4, 1.6
)
footer(sl, 12)

# ══════════════════════════════════════════════════════════════
# SLIDE 13 — GINI → DINERO
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "¿Cómo se Traduce el Gini en Dinero para FinanCrece?",
           "P10, P17, P18 · ROI financiero · Pérdida esperada · Provisiones [SUPUESTO]")

add_image_safe(sl, FIG/"ppt_gini_dinero.png", 0.3, 1.55, 9.8, 4.2)

add_rect(sl, 10.3, 1.55, 5.4, 4.2, BG_CARD)
add_rect(sl, 10.3, 1.55, 0.06, 4.2, GOLD)
add_text(sl, "Cálculo detallado [SUPUESTO]", 10.55, 1.65, 5.0, 0.38,
         size=13, bold=True, color=GOLD)
table_simple(sl,
    ["Escenario", "Resultado (USD)"],
    [
        ["Pérdida sin modelo (aprueba todo)", "-$141,000"],
        ["Ahorro por scoring (modelo activo)", "+$96,600"],
        ["Costo de rechazos erróneos",         "-$12,000 [EST]"],
        ["Beneficio neto estimado",             "+$84,600"],
        ["ROI vs. política base",               "+107%"],
    ],
    10.3, 2.1, 5.4, 3.1
)
add_text(sl, "Gini = 0 → sin discriminación\nGini = 0.665 → ahorro ~$96K\nGini = 1.0 → modelo perfecto",
         10.55, 5.25, 5.1, 0.65, size=10, color=LIGHT, italic=True)

add_rect(sl, 0.3, 5.9, 15.4, 2.6, BG_CARD)
add_rect(sl, 0.3, 5.9, 0.06, 2.6, BLUE)
add_text(sl, "Marco Basilea II: Pérdida Esperada (EL = PD × EAD × LGD)", 0.55, 6.0, 14.8, 0.4,
         size=13, bold=True, color=BLUE)

table_simple(sl,
    ["Banda", "PD (modelo)", "EAD (monto promedio)", "LGD [SUPUESTO]", "EL por cliente", "Acción de provisión"],
    [
        ["Bajo riesgo",   "~8.9%",  "$2,500", "100%", "~$222",  "✅ Provisión mínima"],
        ["Riesgo medio",  "~17.0%", "$2,500", "100%", "~$425",  "⚠️ Provisión media"],
        ["Alto riesgo",   "~59.4%", "$2,500", "100%", "~$1,485","🔴 Provisión alta o rechazar"],
    ],
    0.3, 6.45, 15.4, 1.9
)
footer(sl, 13)

# ══════════════════════════════════════════════════════════════
# SLIDE 14 — POLÍTICA DE 3 BANDAS Y ROI
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Política de 3 Bandas: De la Probabilidad a la Acción",
           "P12 · Herramienta con umbrales ajustables en tiempo real [SUPUESTO]")

add_image_safe(sl, FIG/"ppt_herramienta_ajuste_umbral.png", 0.3, 1.55, 9.5, 4.2)

add_rect(sl, 10.0, 1.55, 5.7, 4.2, BG_CARD)
add_rect(sl, 10.0, 1.55, 0.06, 4.2, ACCENT)
add_text(sl, "Las 3 Bandas (datos reales)", 10.25, 1.65, 5.2, 0.38,
         size=13, bold=True, color=ACCENT)

table_simple(sl,
    ["Banda", "Umbral", "DR real", "N val", "Decisión"],
    [
        ["🟢 BAJO",  "< 0.20",      "8.9%",  "45 (28%)", "Aprobar completo"],
        ["🟡 MEDIO", "0.20 – 0.45", "4.3%",  "46 (29%)", "Condicionar 50%"],
        ["🔴 ALTO",  "> 0.45",      "59.4%", "69 (43%)", "Rechazar"],
    ],
    10.0, 2.1, 5.7, 2.3
)
bullet_block(sl, [
    ("Los umbrales u1=0.20 y u2=0.45 son AJUSTABLES", True, GOLD),
    ("Más conservador → u1=0.15, u2=0.35", False, LIGHT),
    ("Más inclusivo → u1=0.25, u2=0.55", False, LIGHT),
    ("El sistema recalcula el ROI automáticamente", False, ACCENT),
], 10.25, 4.45, 5.2, 1.15, size=10)

add_rect(sl, 0.3, 5.9, 15.4, 2.6, BG_CARD)
add_rect(sl, 0.3, 5.9, 0.06, 2.6, GREEN)
add_text(sl, "Trazabilidad y Transparencia — P13: ¿Cómo explicamos un rechazo?", 0.55, 6.0, 14.8, 0.4,
         size=13, bold=True, color=GREEN)

add_text(sl,
    "Cliente solicita crédito → Sistema calcula prob_default = 0.67 → Asigna banda ALTA → Rechaza",
    0.55, 6.45, 14.8, 0.35, size=11, color=LIGHT, italic=True)

table_simple(sl,
    ["Campo en el reporte de rechazo", "Valor ejemplo"],
    [
        ["prob_default calculada",      "0.67 (Alta: > 0.45)"],
        ["Razón #1 (mayor impacto)",    "checking_status = '<0' → cuenta en negativo"],
        ["Razón #2",                    "cuota_estimada = $485/mes → 41% del ingreso declarado"],
        ["Razón #3",                    "employment = '<1 año' → empleo muy reciente"],
        ["Alternativa ofrecida",        "Línea de $1,000 condicionada con co-deudor [SUPUESTO]"],
    ],
    0.3, 6.85, 15.4, 1.5
)
footer(sl, 14)

# ══════════════════════════════════════════════════════════════
# SLIDE 15 — CONCLUSIONES Y RECOMENDACIONES
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Conclusiones y Recomendaciones para FinanCrece S.A.",
           "P16, P21, P22 · Pricing · Capital · Próximos pasos")

# 3 mensajes principales
for i, (titulo, desc, col, x) in enumerate([
    ("1. SCORING antes del\ndesembolso",
     "LightGBM entrega una prob_default por cliente.\nAUC=0.833 · Gini=0.665 · KS=0.633.\nSuperior a cualquier evaluación subjetiva.",
     ACCENT, 0.3),
    ("2. POLÍTICA DE 3 BANDAS\noperativa",
     "Bajo DR=8.9% → Aprobar completo\nMedio DR=4.3% → Condicionar 50% línea\nAlto DR=59.4% → Rechazar / escalar",
     GOLD, 5.5),
    ("3. ROI MEDIBLE y\nprovisiones optimizadas",
     "+107% ROI vs. aprobar todos [SUPUESTO]\nProvisión diferenciada por banda de riesgo\nLiberación de capital en banda baja",
     GREEN, 10.7),
]):
    add_rect(sl, x, 1.6, 4.9, 3.5, BG_CARD)
    add_rect(sl, x, 1.6, 4.9, 0.06, col)
    add_text(sl, titulo, x+0.15, 1.75, 4.6, 0.85,
             size=14, bold=True, color=col)
    add_text(sl, desc, x+0.15, 2.7, 4.6, 2.1,
             size=11, color=LIGHT)

# Pricing crediticio
add_rect(sl, 0.3, 5.3, 9.8, 3.2, BG_CARD)
add_rect(sl, 0.3, 5.3, 0.06, 3.2, AMBER)
add_text(sl, "Pricing Crediticio por Banda de Riesgo — P21 [SUPUESTO]", 0.55, 5.4, 9.3, 0.38,
         size=13, bold=True, color=AMBER)
add_text(sl,
    "Tasa = Tasa libre de riesgo + Prima PD + Prima LGD + Margen operativo",
    0.55, 5.83, 9.3, 0.35, size=11, color=ACCENT, italic=True)
table_simple(sl,
    ["Banda", "PD estimada", "Prima de riesgo", "Tasa sugerida"],
    [
        ["✅ Bajo riesgo",   "~8.9%",  "+2% pp",  "Tasa base + 2%"],
        ["⚠️ Riesgo medio",  "~17.0%", "+6% pp",  "Tasa base + 6%"],
        ["🔴 Alto riesgo",   "~59.4%", "+15% pp", "Rechazar o garantía"],
    ],
    0.3, 6.2, 9.8, 2.1
)

# Próximos pasos
add_rect(sl, 10.3, 5.3, 5.4, 3.2, BG_CARD)
add_rect(sl, 10.3, 5.3, 0.06, 3.2, BLUE)
add_text(sl, "Próximos pasos", 10.55, 5.4, 5.0, 0.38,
         size=13, bold=True, color=BLUE)
bullet_block(sl, [
    ("Integrar buró crediticio real (SBS / EQUIFAX)", False, LIGHT),
    ("Añadir variables de ingreso verificado y deuda vigente", False, LIGHT),
    ("Monitoreo mensual de PSI (drift de score)", True, ACCENT),
    ("Recalibración semestral del modelo", False, LIGHT),
    ("Sistema de alertas por cohorte de aprobación", False, LIGHT),
    ("Cross-selling a clientes banda baja — P22", True, GOLD),
], 10.55, 5.85, 5.1, 2.5, size=10)

add_rect(sl, 0.3, 8.52, 15.4, 0.18, ACCENT)
add_text(sl,
    '"Un modelo no reemplaza al analista de crédito — le da superpoderes: predecir el riesgo antes de que se materialice"',
    0.3, 8.54, 15.4, 0.22, size=9, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
footer(sl, 15)

# ─── GUARDAR ───────────────────────────────────────────────────
prs.save(OUT)
print(f"\n✅ Presentación guardada en: {OUT}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Tamaño: {OUT.stat().st_size/1024:.0f} KB")
