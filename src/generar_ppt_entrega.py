from pathlib import Path

import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(".")
OUTPUT = PROJECT_ROOT / "GRUPOArmonitech.pptx"
FIG_DIR = PROJECT_ROOT / "reports" / "figures"

NAVY = RGBColor(10, 22, 40)
NAVY_2 = RGBColor(13, 31, 56)
BLUE = RGBColor(45, 109, 181)
GOLD = RGBColor(245, 166, 35)
CYAN = RGBColor(0, 190, 220)
WHITE = RGBColor(250, 252, 255)
MUTED = RGBColor(168, 180, 196)
GREEN = RGBColor(37, 160, 120)
RED = RGBColor(215, 83, 83)


def ensure_extra_figures() -> None:
    df = pd.read_excel(PROJECT_ROOT / "dataInicial" / "dataset_credito-train.xlsx", engine="openpyxl")
    df["default_90d"] = (df["target"] == "bad").astype(int)

    fig, axes = plt.subplots(2, 2, figsize=(12, 7))
    axes = axes.flatten()
    configs = [
        ("purpose", "Riesgo por proposito del credito", 8),
        ("employment", "Riesgo por estabilidad laboral", 8),
        ("housing", "Riesgo por tipo de vivienda", 8),
        ("personal_status", "Riesgo por estado personal codificado", 8),
    ]
    for ax, (col, title, limit) in zip(axes, configs):
        rates = df.groupby(col)["default_90d"].agg(["mean", "count"]).sort_values("mean", ascending=False)
        rates = rates.head(limit).sort_values("mean")
        ax.barh(rates.index.astype(str), rates["mean"], color=[(245 / 255, 166 / 255, 35 / 255)])
        ax.axvline(df["default_90d"].mean(), color=(45 / 255, 109 / 255, 181 / 255), linestyle="--", linewidth=1.5)
        ax.set_title(title, fontsize=11, weight="bold")
        ax.set_xlabel("Default rate")
        ax.xaxis.set_major_formatter(lambda x, pos: f"{x:.0%}")
        ax.tick_params(axis="y", labelsize=8)
    plt.tight_layout()
    out = FIG_DIR / "ppt_09_segmentos_eda_preguntas.png"
    plt.savefig(out, dpi=180, bbox_inches="tight")
    plt.close(fig)

    fig, ax = plt.subplots(figsize=(9, 5))
    scatter = ax.scatter(
        df["duration"],
        df["credit_amount"],
        c=df["default_90d"],
        cmap="coolwarm",
        alpha=0.70,
        edgecolors="white",
        linewidths=0.3,
    )
    ax.set_title("Monto solicitado vs duracion del credito", fontsize=13, weight="bold")
    ax.set_xlabel("Duracion del prestamo")
    ax.set_ylabel("Monto solicitado")
    ax.grid(alpha=0.25)
    cbar = plt.colorbar(scatter, ax=ax)
    cbar.set_label("0=good, 1=bad")
    plt.tight_layout()
    out = FIG_DIR / "ppt_10_monto_duracion.png"
    plt.savefig(out, dpi=180, bbox_inches="tight")
    plt.close(fig)


def add_bg(slide, color=NAVY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_text(slide, text, x, y, w, h, size=20, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, eyebrow=None):
    if eyebrow:
        add_text(slide, eyebrow.upper(), 0.65, 0.28, 7.5, 0.28, 9, CYAN, True)
    add_text(slide, title, 0.65, 0.55, 8.6, 0.55, 24, WHITE, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.18), Inches(1.3), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()


def add_footer(slide, n):
    add_text(slide, "Datathon ESAN 2026 | Grupo Armonitech", 0.65, 7.05, 5.5, 0.25, 8, MUTED)
    add_text(slide, f"{n:02d}", 12.35, 7.04, 0.35, 0.25, 8, MUTED, align=PP_ALIGN.RIGHT)


def add_bullets(slide, bullets, x, y, w, h, size=13, color=WHITE, gap=0.13):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    for idx, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(gap * 20)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def add_metric(slide, value, label, x, y, w=1.65, accent=GOLD):
    add_text(slide, value, x, y, w, 0.42, 22, accent, True, PP_ALIGN.CENTER)
    add_text(slide, label, x, y + 0.45, w, 0.38, 8.5, MUTED, False, PP_ALIGN.CENTER)


def add_image(slide, path, x, y, w=None, h=None):
    path = Path(path)
    if not path.exists():
        return None
    if w and h:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if w:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    if h:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), height=Inches(h))
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def add_tag(slide, text, x, y, w, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.33))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text(slide, text, x + 0.05, y + 0.055, w - 0.1, 0.18, 8.5, WHITE, True, PP_ALIGN.CENTER)


def build_ppt():
    ensure_extra_figures()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 Cover
    slide = prs.slides.add_slide(blank)
    add_bg(slide, NAVY)
    add_text(slide, "FinanCrece S.A.", 0.8, 0.72, 5.4, 0.5, 16, CYAN, True)
    add_text(slide, "Scoring de riesgo crediticio", 0.75, 1.35, 8.8, 0.78, 35, WHITE, True)
    add_text(slide, "Probabilidad de default, politica de aprobacion y herramienta operativa para originacion.", 0.78, 2.25, 8.4, 0.55, 16, MUTED)
    add_metric(slide, "0.802", "AUC validacion", 0.85, 4.1)
    add_metric(slide, "0.604", "Gini", 2.65, 4.1, accent=CYAN)
    add_metric(slide, "USD 99.6K", "ahorro estimado", 4.45, 4.1, w=2.0)
    add_metric(slide, "3 bandas", "decision de riesgo", 6.75, 4.1, w=1.8, accent=GREEN)
    add_image(slide, FIG_DIR / "ppt_07_distribucion_scores_test.png", 8.65, 1.45, w=3.85)
    add_footer(slide, 1)

    # 2 Business objectives
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Objetivos del negocio", "1. Objetivos")
    add_text(slide, "Reto", 0.75, 1.6, 2.2, 0.3, 14, GOLD, True)
    add_bullets(slide, [
        "Estimar probabilidad de incumplimiento antes del desembolso.",
        "Aprobar buenos clientes sin castigar inclusion financiera.",
        "Reducir perdida esperada, provisiones mal asignadas y capital en riesgo.",
        "Convertir el score en una politica simple: aprobar, condicionar o rechazar.",
    ], 0.8, 2.05, 5.1, 2.25, 13)
    add_text(slide, "Respuesta al requerimiento oficial", 6.55, 1.6, 4.8, 0.35, 14, GOLD, True)
    add_bullets(slide, [
        "Entrada: caracteristicas del solicitante.",
        "Proceso: pipeline predictivo anti-leakage.",
        "Salida: prob_default + banda de riesgo + decision.",
        "Uso: herramienta de scoring individual y batch.",
    ], 6.6, 2.05, 5.5, 1.8, 13)
    add_tag(slide, "Inclusion responsable", 6.65, 4.45, 1.85, GREEN)
    add_tag(slide, "Control de riesgo", 8.75, 4.45, 1.75, BLUE)
    add_tag(slide, "Rentabilidad", 10.72, 4.45, 1.35, GOLD)
    add_text(slide, "Idea clave: el modelo no reemplaza la politica; la vuelve medible, trazable y ajustable.", 1.05, 5.7, 10.8, 0.45, 18, WHITE, True, PP_ALIGN.CENTER)
    add_footer(slide, 2)

    # 3 Data processing
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Procesamiento de datos", "2. Procesamiento")
    add_bullets(slide, [
        "Train oficial: 800 clientes, 22 columnas; test oficial: 200 clientes, 23 columnas.",
        "Target raw: target = good/bad; target canonico: default_90d, donde bad = 1.",
        "No se inventaron variables: el Excel no contiene buro externo, ingreso, geografia ni canal.",
        "Feature engineering fila a fila: carga financiera, cuota estimada, ordinales de riesgo e interacciones.",
        "Preprocesamiento dentro del pipeline: imputacion, escalado y one-hot ajustados solo con train.",
    ], 0.8, 1.55, 6.0, 4.2, 12.5)
    add_image(slide, FIG_DIR / "notebook_eda_target_checking.png", 7.1, 1.45, w=5.25)
    add_footer(slide, 3)

    # 4 EDA - why credit and purposes
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Analisis exploratorio: proposito, empleo, vivienda y estado personal", "3. EDA")
    add_image(slide, FIG_DIR / "ppt_09_segmentos_eda_preguntas.png", 0.65, 1.35, w=7.2)
    add_text(slide, "Lecturas para defensa", 8.25, 1.45, 3.4, 0.3, 14, GOLD, True)
    add_bullets(slide, [
        "La data observa el proposito del credito, no motivaciones psicologicas: consumo, auto, negocio, educacion y reparaciones.",
        "Empleo estable reduce riesgo; los grupos de 4-7 y >=7 anos muestran mejor comportamiento que <1 ano.",
        "Vivienda propia tiende a menor default que alquiler o vivienda gratuita.",
        "Estado personal influye menos que liquidez, historial y ahorros; se interpreta con prudencia regulatoria.",
    ], 8.25, 1.95, 4.4, 3.7, 12)
    add_footer(slide, 4)

    # 5 EDA risk drivers
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Que explica el incumplimiento", "3. EDA")
    add_text(slide, "Principales senales observadas", 0.8, 1.45, 3.8, 0.3, 14, GOLD, True)
    add_bullets(slide, [
        "Liquidez debil: checking_status <0 eleva el default.",
        "Historial riesgoso: no credits/all paid y all paid concentran alto default.",
        "Ahorros bajos aumentan vulnerabilidad; el efecto no es perfectamente lineal porque 'no known savings' es una categoria informativa.",
        "Monto y duracion importan juntos: la carga financiera captura exposicion acumulada.",
        "Installment_commitment aproxima presion de cuota sobre capacidad de pago.",
    ], 0.8, 1.95, 5.4, 4.25, 12.2)
    add_image(slide, FIG_DIR / "ppt_10_monto_duracion.png", 6.7, 1.35, w=5.6)
    add_footer(slide, 5)

    # 6 Feature engineering + interpretability
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Feature engineering e interpretabilidad", "4. Desarrollo del modelo")
    add_image(slide, FIG_DIR / "ppt_03_interpretabilidad_coeficientes.png", 0.65, 1.35, w=6.4)
    add_text(slide, "Trazabilidad ante reclamos", 7.35, 1.45, 4.6, 0.3, 14, GOLD, True)
    add_bullets(slide, [
        "Cada decision se explica por variables observables: liquidez, historial, ahorro, empleo, monto y plazo.",
        "La herramienta devuelve prob_default, banda y factores de alerta.",
        "Los limites se ajustan cambiando u_bajo, u_alto o el threshold financiero.",
        "Riesgos ocultos no visibles se mitigan con proxies y se recomiendan como futuras fuentes: buro real, ingresos, deuda externa y macro.",
    ], 7.35, 1.95, 4.9, 3.8, 12)
    add_footer(slide, 6)

    # 7 Model development
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Desarrollo y seleccion del modelo", "4. Desarrollo del modelo")
    add_image(slide, FIG_DIR / "ppt_02_comparacion_modelos.png", 0.65, 1.35, w=6.4)
    add_text(slide, "Campeon operativo", 7.35, 1.4, 4.2, 0.3, 14, GOLD, True)
    add_bullets(slide, [
        "Modelo seleccionado: LogReg_balanced_C01.",
        "No elegimos solo el AUC bruto: penalizamos sobreajuste y mala calibracion.",
        "La regresion logistica regularizada es interpretable, portable y adecuada para explicar aprobaciones.",
        "Split estratificado por ausencia de columna temporal oficial.",
    ], 7.35, 1.85, 4.8, 2.4, 12.5)
    add_metric(slide, "0.802", "AUC val", 7.35, 4.75, 1.3)
    add_metric(slide, "0.604", "Gini", 8.8, 4.75, 1.2, CYAN)
    add_metric(slide, "0.527", "KS", 10.05, 4.75, 1.2, GREEN)
    add_metric(slide, "0.170", "Brier", 11.25, 4.75, 1.2)
    add_footer(slide, 7)

    # 8 Evaluation
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Evaluacion tecnica del score", "4. Evaluacion")
    add_image(slide, FIG_DIR / "ppt_04_curvas_desempeno.png", 0.7, 1.35, w=6.7)
    add_text(slide, "Lectura de riesgo", 7.7, 1.5, 4.2, 0.3, 14, GOLD, True)
    add_bullets(slide, [
        "Gini alto significa mejor ordenamiento de riesgo; monetizamos ese ordenamiento con ROI por threshold.",
        "Brier se revisa porque la salida se usa como probabilidad, no solo como ranking.",
        "El test interno cae a AUC 0.705: se reporta con transparencia y se compensa con monitoreo y recalibracion.",
        "En crisis economicas: endurecer umbrales, subir revision manual y recalibrar con mora reciente.",
    ], 7.7, 2.0, 4.7, 3.6, 12)
    add_footer(slide, 8)

    # 9 Policy & financial impact
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Politica de riesgo e impacto financiero", "5. Negocio")
    add_image(slide, FIG_DIR / "notebook_politica_bandas.png", 0.65, 1.35, w=4.7)
    add_image(slide, FIG_DIR / "ppt_05_roi_threshold.png", 5.25, 1.35, w=4.4)
    add_text(slide, "Impacto", 10.05, 1.45, 2.5, 0.3, 14, GOLD, True)
    add_bullets(slide, [
        "Threshold optimo: 0.36.",
        "Ahorro incremental estimado: USD 99,600.",
        "ROI vs base: +110.48%.",
        "PD alimenta perdida esperada, provisiones, pricing e intensidad de revision.",
        "Mejor asignacion de capital: mas capital a buenos riesgos, menor exposicion a defaults esperados.",
    ], 10.05, 1.95, 2.7, 4.0, 11)
    add_footer(slide, 9)

    # 10 Conclusions and tool
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Conclusiones y recomendacion al jurado", "6. Conclusiones")
    add_image(slide, FIG_DIR / "ppt_08_ejemplo_herramienta_scoring.png", 0.65, 1.35, w=6.1)
    add_text(slide, "Cierre ejecutivo", 7.15, 1.45, 4.6, 0.3, 14, GOLD, True)
    add_bullets(slide, [
        "La herramienta cumple el requerimiento: ingresa caracteristicas, retorna prob_default y decision.",
        "La politica equilibra inclusion financiera y control de riesgo: aprueba bajo riesgo, condiciona riesgo medio y mitiga alto riesgo.",
        "El modelo permite anticipar perdidas, ajustar provisiones y proteger capital.",
        "Siguiente etapa: integrar buro, ingresos, deuda externa, variables macro y monitoreo mensual de drift.",
    ], 7.15, 1.95, 5.2, 3.25, 12)
    add_text(slide, "Mensaje final: scoring interpretable + politica accionable + ROI medible.", 7.2, 5.85, 4.7, 0.55, 16, WHITE, True)
    add_footer(slide, 10)

    prs.save(OUTPUT)


if __name__ == "__main__":
    build_ppt()
    print(f"PowerPoint generado: {OUTPUT}")
